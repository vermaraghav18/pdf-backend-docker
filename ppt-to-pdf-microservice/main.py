from fastapi import FastAPI, UploadFile
from fastapi.responses import FileResponse
import os
import fitz  # PyMuPDF
import tempfile
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

app = FastAPI()

@app.post("/")
async def convert_ppt_to_pdf(file: UploadFile):
    temp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(temp_dir, file.filename)

    with open(pptx_path, "wb") as f:
        f.write(await file.read())

    prs = Presentation(pptx_path)
    pdf_path = pptx_path.replace(".pptx", ".pdf")

    c = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in prs.slides:
        text = ""
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
        c.setFont("Helvetica", 12)
        y = height - 50
        for line in text.splitlines():
            c.drawString(50, y, line)
            y -= 14
            if y < 40:
                c.showPage()
                y = height - 50
        c.showPage()

    c.save()

    return FileResponse(pdf_path, filename="converted.pdf")
